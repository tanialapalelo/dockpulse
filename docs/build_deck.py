"""Generate the DockPulse pitch deck (.pptx) from the real benchmark numbers.

Run:  python docs/build_deck.py
Output: docs/DockPulse_Deck.pptx  (10 slides, 16:9)

Re-run any time data/benchmark.json changes and the numbers update automatically.
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
BENCH = json.loads((ROOT / "data" / "benchmark.json").read_text())

APP_URL = "https://dockpulse-qdswtbn7xx5t6yje6ksyuy.streamlit.app"
REPO_URL = "https://github.com/tanialapalelo/dockpulse"

# ---- palette --------------------------------------------------------------
NVIDIA = RGBColor(0x76, 0xB9, 0x00)
GOOGLE = RGBColor(0x42, 0x85, 0xF4)
INK = RGBColor(0x11, 0x14, 0x1A)
GREY = RGBColor(0x5F, 0x6B, 0x7A)
LIGHT = RGBColor(0xF3, 0xF5, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xD6, 0x28, 0x28)
BLUE = RGBColor(0x1D, 0x6F, 0xD6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

N_TRIPS = f"{BENCH['n_trips']:,}"
SPEEDUP = f"{BENCH['speedup']}×"
CPU_S = f"{BENCH['cpu_seconds']}s"
GPU_S = f"{BENCH['gpu_seconds']}s"
GPU_NAME = BENCH["gpu_name"]


def _set(fr, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT, font="Calibri"):
    fr.word_wrap = True
    p = fr.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return p


def box(slide, x, y, w, h, text, size, color, **kw):
    tb = slide.shapes.add_textbox(x, y, w, h)
    _set(tb.text_frame, text, size, color, **kw)
    return tb


def bullets(slide, x, y, w, h, items, size=18, color=INK, gap=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, col, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.color.rgb = col
        r.font.bold = bold
        r.font.name = "Calibri"
    return tb


def fill(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def bar(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def chip(slide, x, y, text, color=NVIDIA):
    from pptx.enum.shapes import MSO_SHAPE
    w = Inches(0.18 + 0.11 * len(text))
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.42))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = False
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    _set(tf, text, 11, WHITE, bold=True, align=PP_ALIGN.CENTER)
    return x + w + Inches(0.12)


def header(slide, kicker, title, accent=NVIDIA):
    bar(slide, Inches(0.6), Inches(0.55), Inches(0.12), Inches(0.9), accent)
    box(slide, Inches(0.85), Inches(0.5), Inches(11.8), Inches(0.4),
        kicker, 13, accent, bold=True)
    box(slide, Inches(0.85), Inches(0.82), Inches(11.8), Inches(0.9),
        title, 30, INK, bold=True)


def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------- Slide 1
s = prs.slides.add_slide(BLANK)
fill(s, INK)
bar(s, 0, Inches(3.05), SW, Inches(0.06), NVIDIA)
box(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.4),
    "DockPulse", 66, WHITE, bold=True)
box(s, Inches(0.95), Inches(2.75), Inches(11.5), Inches(0.6),
    "Will your docks run empty next hour?", 26, NVIDIA, bold=True)
box(s, Inches(0.95), Inches(3.3), Inches(11.5), Inches(0.9),
    "GPU-accelerated rebalancing early-warning for bike-share operations.",
    18, RGBColor(0xC7, 0xCE, 0xD8))
cx = Inches(0.95)
for t, c in [("BigQuery", GOOGLE), ("RAPIDS cuDF", NVIDIA), ("Gemini", GOOGLE), ("Streamlit", GREY)]:
    cx = chip(s, cx, Inches(4.35), t, c)
box(s, Inches(0.95), Inches(6.4), Inches(11.5), Inches(0.5),
    "tania  ·  Hack2Skill × Google Cloud × NVIDIA — GenAI Academy APAC Edition",
    14, RGBColor(0x9A, 0xA4, 0xB2))
note(s, "Every bike-share dispatcher fights the same hourly battle — I built an early-warning board for it, accelerated on a GPU.")

# ---------------------------------------------------------------- Slide 2
s = prs.slides.add_slide(BLANK); fill(s, WHITE)
header(s, "THE REAL USER & THE STUCK DECISION", "A dispatcher makes the same call every hour")
bullets(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(4), [
    ("WHO", NVIDIA, True),
    ("A bike-share operations dispatcher — and every commuter who's found an empty dock.", INK, False),
    ("", INK, False),
    ("THE DECISION, EVERY HOUR", NVIDIA, True),
    ("Where do I send rebalancing crews before stations run empty or overflow?", INK, False),
    ("", INK, False),
    ("TODAY", RED, True),
    ("Reactive. Crews arrive AFTER riders are already stranded at an empty dock.", INK, False),
], size=20, gap=6)
note(s, "Name the person and the recurring, data-dependent decision. Judges reward a specific real user.")

# ---------------------------------------------------------------- Slide 3
s = prs.slides.add_slide(BLANK); fill(s, WHITE)
header(s, "WHY IT'S HARD", "The bottleneck is data volume vs. time-to-decision")
bullets(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(3.5), [
    (f"•  {N_TRIPS} real trips must become per-station, per-hour flow — every refresh.", INK, False),
    ("•  On CPU that aggregation is slow, so the insight arrives too late to act on.", INK, False),
    ("•  Stale forecast  →  empty dock  →  lost rider.", RED, True),
], size=21, gap=14)
box(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1),
    "The decision doesn't need more data — it needs the SAME data, faster.", 20, GOOGLE, bold=True, italic=True)
note(s, "This sets up why acceleration matters: the blocker is compute time, not missing data.")

# ---------------------------------------------------------------- Slide 4
s = prs.slides.add_slide(BLANK); fill(s, WHITE)
header(s, "THE SOLUTION", "A live risk board the dispatcher acts on")
bullets(s, Inches(0.9), Inches(2.0), Inches(6.0), Inches(4.5), [
    ("\U0001F5FA️  Risk map", INK, True),
    ("Every station coloured by empty/full risk for the selected hour.", GREY, False),
    ("\U0001F6A8  Ranked 'rebalance now' alerts", INK, True),
    ("Top stations to act on + the concrete action (move ~N bikes).", GREY, False),
    ("\U0001F9E0  Gemini ops briefing", INK, True),
    ("The top alerts summarised into a 20-second read.", GREY, False),
    ("⚡  Acceleration proof", INK, True),
    ("CPU vs GPU timing, in the app.", GREY, False),
], size=17, gap=4)
bar(s, Inches(7.2), Inches(2.0), Inches(5.2), Inches(4.3), LIGHT)
box(s, Inches(7.2), Inches(3.7), Inches(5.2), Inches(0.9),
    "[ Paste app screenshot here ]", 15, GREY, align=PP_ALIGN.CENTER)
note(s, "Show, don't tell. Point at the top alert and read the action out loud.")

# ---------------------------------------------------------------- Slide 5
s = prs.slides.add_slide(BLANK); fill(s, WHITE)
header(s, "ARCHITECTURE", "Two clean layers: GPU offline, app lightweight", GOOGLE)
steps = [
    ("Google Cloud data", "BigQuery\nlondon_bicycles", GOOGLE),
    ("NVIDIA acceleration", "Colab T4\npandas vs cuDF", NVIDIA),
    ("Artifacts", "features.parquet\nbenchmark.json", GREY),
    ("Decision layer", "Streamlit app\n+ Gemini briefing", GOOGLE),
]
x = Inches(0.9)
for i, (cap, body, col) in enumerate(steps):
    bar(s, x, Inches(2.6), Inches(2.7), Inches(1.7), col)
    tb = s.shapes.add_textbox(x, Inches(2.6), Inches(2.7), Inches(1.7))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set(tf, body, 15, WHITE, bold=True, align=PP_ALIGN.CENTER)
    box(s, x, Inches(4.4), Inches(2.7), Inches(0.5), cap, 12, col, bold=True, align=PP_ALIGN.CENTER)
    if i < 3:
        box(s, x + Inches(2.62), Inches(2.9), Inches(0.4), Inches(1), "→", 28, INK, align=PP_ALIGN.CENTER)
    x = Emu(int(x) + int(Inches(3.05)))
box(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.4),
    "Heavy GPU work runs ONCE offline and is benchmarked; the app just serves tiny files — "
    "so it deploys free, needs no GPU at runtime, and never breaks during the demo.",
    16, INK, italic=True)
note(s, "Explain the split: reproducible GPU proof + a lightweight, always-up app.")

# ---------------------------------------------------------------- Slide 6
s = prs.slides.add_slide(BLANK); fill(s, WHITE)
header(s, "THE DATA PIPELINE", "Traceable from ingest to action", GOOGLE)
pipe = [
    ("1  Ingest + Clean", f"{N_TRIPS} London trips from BigQuery; attach station geo & capacity.", GOOGLE),
    ("2  Analyze + Model", "Net flow per station × hour-of-week + volatility (GPU cuDF).", NVIDIA),
    ("3  Score", "Empty/full risk 0–100 + bikes-to-move per station.", INK),
    ("4  Visualize + Act", "Risk map, ranked alerts, Gemini briefing.", GOOGLE),
]
y = Inches(2.2)
for cap, body, col in pipe:
    bar(s, Inches(0.9), y, Inches(0.12), Inches(0.9), col)
    box(s, Inches(1.15), y, Inches(3.4), Inches(0.9), cap, 18, col, bold=True)
    box(s, Inches(4.6), y - Emu(int(Pt(2))), Inches(8), Inches(0.9), body, 16, INK)
    y = Emu(int(y) + int(Inches(1.15)))
note(s, "Mirror NVIDIA's 'Required Data Pipeline' slide — trace it end to end.")

# ---------------------------------------------------------------- Slide 7 (money)
s = prs.slides.add_slide(BLANK); fill(s, INK)
bar(s, Inches(0.6), Inches(0.55), Inches(0.12), Inches(0.9), NVIDIA)
box(s, Inches(0.85), Inches(0.5), Inches(11.8), Inches(0.4), "ACCELERATION PROOF", 13, NVIDIA, bold=True)
box(s, Inches(0.85), Inches(0.82), Inches(11.8), Inches(0.9),
    "Same code. Swap CPU → GPU.", 30, WHITE, bold=True)
box(s, Inches(0.9), Inches(2.3), Inches(5.6), Inches(2), "pandas (CPU)", 22, RGBColor(0xC7,0xCE,0xD8))
box(s, Inches(0.9), Inches(2.8), Inches(5.6), Inches(1.2), CPU_S, 60, RED, bold=True)
box(s, Inches(0.9), Inches(4.4), Inches(5.6), Inches(0.6), "cuDF (GPU)", 22, RGBColor(0xC7,0xCE,0xD8))
box(s, Inches(0.9), Inches(4.9), Inches(5.6), Inches(1.2), GPU_S, 60, NVIDIA, bold=True)
bar(s, Inches(6.9), Inches(2.3), Inches(5.5), Inches(3.8), RGBColor(0x1A,0x24,0x12))
box(s, Inches(6.9), Inches(2.8), Inches(5.5), Inches(1), SPEEDUP, 96, NVIDIA, bold=True, align=PP_ALIGN.CENTER)
box(s, Inches(6.9), Inches(4.5), Inches(5.5), Inches(0.6), "faster time-to-insight", 20, WHITE, align=PP_ALIGN.CENTER)
box(s, Inches(6.9), Inches(5.15), Inches(5.5), Inches(0.7),
    f"{N_TRIPS} trips  ·  {GPU_NAME}", 14, RGBColor(0x9A,0xA4,0xB2), align=PP_ALIGN.CENTER)
box(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.7),
    "Minutes of waiting become seconds — so the dispatcher acts while the insight is still useful.",
    16, RGBColor(0xC7,0xCE,0xD8), italic=True)
note(s, f"The slide the NVIDIA judge waits for. Say '{SPEEDUP} faster' twice. Real numbers on {GPU_NAME}.")

# ---------------------------------------------------------------- Slide 8
s = prs.slides.add_slide(BLANK); fill(s, WHITE)
header(s, "GEMINI IN THE LOOP", "From risk scores to a plain-English decision")
bullets(s, Inches(0.9), Inches(2.1), Inches(6.0), Inches(4), [
    ("Gemini turns the top alerts into a confident shift briefing:", INK, False),
    ("", INK, False),
    ("•  Names the 2–3 most urgent stations", INK, False),
    ("•  Says whether to send bikes IN or free docks", INK, False),
    ("•  Gives one clear next action", INK, False),
    ("", INK, False),
    ("Template fallback if the API ever fails — the demo never breaks.", GOOGLE, True),
], size=18, gap=6)
bar(s, Inches(7.2), Inches(2.0), Inches(5.2), Inches(4.3), LIGHT)
box(s, Inches(7.2), Inches(3.7), Inches(5.2), Inches(0.9),
    "[ Paste Gemini briefing screenshot ]", 15, GREY, align=PP_ALIGN.CENTER)
note(s, "Shows GenAI doing real decision-support work, not bolted on.")

# ---------------------------------------------------------------- Slide 9
s = prs.slides.add_slide(BLANK); fill(s, WHITE)
header(s, "IMPACT & WHAT'S NEXT", "Act before failure, not after")
bullets(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(4), [
    ("IMPACT", NVIDIA, True),
    ("Fewer empty docks, shorter rider waits, crews dispatched before a station fails.", INK, False),
    ("Validated on real London data: Waterloo Station flagged FULL at evening rush; Soho & Holborn flagged EMPTY.", GREY, False),
    ("", INK, False),
    ("SCALES", GOOGLE, True),
    ("Any city with an open trip feed. Freshness improves with GPU — re-score every hour.", INK, False),
    ("", INK, False),
    ("NEXT", INK, True),
    ("Live streaming feed · XGBoost forecast + weather · Cloud Run hourly re-score.", GREY, False),
], size=17, gap=5)
note(s, "Connect back to 'AI for better living / smarter communities'.")

# ---------------------------------------------------------------- Slide 10
s = prs.slides.add_slide(BLANK); fill(s, INK)
bar(s, 0, Inches(2.9), SW, Inches(0.06), NVIDIA)
box(s, Inches(0.9), Inches(1.2), Inches(11.5), Inches(1), "Try it.", 48, WHITE, bold=True)
bullets(s, Inches(0.9), Inches(3.3), Inches(11.5), Inches(2.5), [
    (f"Live app:      {APP_URL}", NVIDIA, True),
    (f"GitHub:        {REPO_URL}", RGBColor(0xC7,0xCE,0xD8), False),
    ("Demo video:  [ paste your video link ]", RGBColor(0xC7,0xCE,0xD8), False),
], size=18, gap=12)
cx = Inches(0.95)
for t, c in [("BigQuery", GOOGLE), ("RAPIDS cuDF", NVIDIA), ("Gemini", GOOGLE), ("Streamlit", GREY)]:
    cx = chip(s, cx, Inches(6.2), t, c)
box(s, Inches(9.0), Inches(6.7), Inches(3.8), Inches(0.5), "Built in 24 hours.", 16, NVIDIA, bold=True, align=PP_ALIGN.RIGHT)
note(s, "End on the working link. Invite them to click a station.")

# ---------------------------------------------------------------- save
OUT = ROOT / "docs" / "DockPulse_Deck.pptx"
prs.save(OUT)
print(f"Saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
print(f"Numbers baked in: {N_TRIPS} trips, {CPU_S} CPU vs {GPU_S} GPU = {SPEEDUP} on {GPU_NAME}")
